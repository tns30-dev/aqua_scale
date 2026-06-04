"""
Generate AquaShield_Presentation.pptx from 1.md … 16.md in this folder.

Each slide:
  - Title (16:9, white background)
  - Body bullets / table summary (from "### Body" sections)
  - Diagram placeholder ("📌 INSERT DIAGRAM" box) wherever a ```mermaid block
    or "[INSERT … here]" marker appears in the MD
  - Speaker notes pane gets the full "## Speaker notes" section

Run with anaconda Python (which has python-pptx 1.0.2):
    /opt/anaconda3/bin/python3 generate_pptx.py
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

ROOT = Path(__file__).parent
OUT = ROOT / "AquaShield_Presentation.pptx"

# 16 content slides + 1 title slide
SLIDE_FILES = [ROOT / f"{i}.md" for i in range(1, 17)]


# ---------------------------------------------------------------------------
# Markdown parsing helpers
# ---------------------------------------------------------------------------

def read(path: Path) -> str:
    return path.read_text() if path.exists() else ""


def strip_emoji_prefix(s: str) -> str:
    # Trim a leading emoji-and-space pattern, but keep the rest.
    return s


def extract_section(md: str, heading: str) -> str:
    """Extract text under a markdown heading (e.g. '## Speaker notes')
    up to the next heading of the same level or higher."""
    level = heading.count("#")
    pattern = re.compile(
        rf"^{re.escape(heading)}\b(.*?)(?=^{'#' * level}\s|\Z)",
        re.MULTILINE | re.DOTALL,
    )
    m = pattern.search(md)
    if not m:
        return ""
    return m.group(1).strip()


def extract_title(md: str, filename: str) -> str:
    # First, look for the H1 "# (Slide|Step) N — Title"
    m = re.search(r"^#\s+(?:Slide|Step)\s+\d+\s+[—-]+\s+(.+)$", md, re.MULTILINE)
    if m:
        return m.group(1).strip()
    # Then any H1
    m = re.search(r"^#\s+(.+)$", md, re.MULTILINE)
    if m:
        return m.group(1).strip()
    return filename


def extract_slide_subtitle(md: str) -> str:
    # First, look for "### Title" body content (the bolded text after it)
    title_section = re.search(r"^###\s+Title\s*$\n+(.+?)(?=\n*###|\Z)", md, re.MULTILINE | re.DOTALL)
    if title_section:
        text = title_section.group(1).strip()
        # Strip outer ** if present
        text = re.sub(r"^\*\*(.+?)\*\*$", r"\1", text.strip())
        return text
    return ""


def remove_code_blocks(md: str) -> tuple[str, int]:
    """Remove fenced code blocks (mermaid etc); return (cleaned_text, count)."""
    pattern = re.compile(r"^```.*?^```", re.MULTILINE | re.DOTALL)
    count = len(pattern.findall(md))
    cleaned = pattern.sub("", md)
    return cleaned, count


def extract_insert_markers(md: str) -> list[str]:
    """Find all '[INSERT ... here]' or '[INSERT ... ]' markers in the MD."""
    return re.findall(r"\[INSERT[^\]]+\]", md)


def extract_body_bullets(md: str) -> list[str]:
    """
    Pull bullet-friendly content from the 'Slide content' region.
    We:
      - Take everything between '## Slide content' and the next '## Speaker notes'
        (or end of file).
      - Strip mermaid code fences entirely (those are diagrams).
      - Convert tables into a single summary line "  · <columns>".
      - Keep markdown list items as-is.
      - Skip image placeholder lines (already captured separately).
    """
    body_section = extract_section(md, "## Slide content (what appears on the slide)")
    if not body_section:
        # Fallback: everything between the H1 line and the first "## Speaker notes"
        # (or end of file). Strips the front-matter (Time budget · Audience · Framing).
        after_h1 = re.sub(r"^#\s.+?\n", "", md, count=1, flags=re.MULTILINE)
        before_speakers = re.split(r"^##\s+Speaker notes", after_h1, flags=re.MULTILINE)[0]
        body_section = before_speakers

    # Remove mermaid code fences
    body_section, _ = remove_code_blocks(body_section)

    lines = body_section.splitlines()
    bullets: list[str] = []
    in_table = False
    table_header: list[str] = []
    table_rows = 0

    for raw in lines:
        line = raw.rstrip()

        if not line.strip():
            if in_table:
                # flush table
                if table_header:
                    bullets.append(
                        f"📊 Table: {' · '.join(table_header)}  (rows: {table_rows})"
                    )
                in_table = False
                table_header = []
                table_rows = 0
            continue

        # markdown headers inside body (### Subtitle etc.) → make bold lines
        if line.startswith("####"):
            txt = line.lstrip("#").strip()
            bullets.append(f"▸ {txt}")
            continue
        if line.startswith("###"):
            txt = line.lstrip("#").strip()
            bullets.append(f"■ {txt}")
            continue

        # tables
        if line.startswith("|"):
            cells = [c.strip() for c in line.strip("|").split("|")]
            if not in_table:
                # header row
                in_table = True
                table_header = cells
                table_rows = 0
            elif re.match(r"^\|[\s\-:|]+\|?$", line.strip()):
                continue  # separator row
            else:
                table_rows += 1
            continue

        # quote lines, like "> [INSERT ...]"
        if line.startswith(">"):
            txt = line.lstrip(">").strip()
            if "INSERT" in txt:
                bullets.append(f"🖼️  {txt}")
            elif txt:
                bullets.append(f"  ▹ {txt}")
            continue

        # bullet lines
        if line.lstrip().startswith(("- ", "* ", "• ")):
            bullets.append(line)
            continue

        # other meaningful prose lines (skip notes-style)
        if line.strip().startswith("*Last updated"):
            continue
        if line.strip() == "---":
            continue
        # generic text — make a bullet
        if len(line.strip()) > 0 and not line.startswith("#"):
            bullets.append(line.strip())

    # flush trailing table
    if in_table and table_header:
        bullets.append(
            f"📊 Table: {' · '.join(table_header)}  (rows: {table_rows})"
        )

    return bullets


def extract_speaker_notes(md: str) -> str:
    section = extract_section(md, "## Speaker notes")
    if not section:
        return ""
    # Strip the leading "(≈ 60s — say out loud …)" parenthetical-only line
    section = re.sub(r"^\(.*?\)\s*\n", "", section, flags=re.MULTILINE)
    # Convert blockquote ">" prefixes to plain text
    lines = [re.sub(r"^>\s?", "", ln) for ln in section.splitlines()]
    return "\n".join(lines).strip()


def count_mermaid_blocks(md: str) -> int:
    return len(re.findall(r"^```mermaid", md, re.MULTILINE))


# ---------------------------------------------------------------------------
# PPTX builder
# ---------------------------------------------------------------------------

# AquaShield palette
TEAL = RGBColor(0x00, 0x76, 0x7B)
DEEP = RGBColor(0x00, 0x3D, 0x4C)
SAND = RGBColor(0xF5, 0xF1, 0xE8)
SLATE = RGBColor(0x37, 0x47, 0x4F)
ACCENT = RGBColor(0xE6, 0x7E, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def make_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill_slide_bg(slide, DEEP)

    # Title text box
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12.0), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "AquaShield Monitoring v2"
    r.font.size = Pt(60)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(12.0), Inches(1.0))
    sf = sub.text_frame
    sf.word_wrap = True
    sp = sf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = "Helping aquaculture farms catch water-quality problems before they cost a harvest"
    sr.font.size = Pt(20)
    sr.font.italic = True
    sr.font.color.rgb = SAND

    # Footer line
    foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.5))
    ff = foot.text_frame
    fp = ff.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = "MTech in Software Engineering · Final Presentation · Thet Naung Soe"
    fr.font.size = Pt(14)
    fr.font.color.rgb = SAND


def _fill_slide_bg(slide, rgb):
    from pptx.oxml.ns import qn

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb


def make_content_slide(prs: Presentation, slide_num: int, md_path: Path) -> None:
    md = read(md_path)
    if not md:
        return

    title_text = extract_title(md, md_path.stem)
    subtitle = extract_slide_subtitle(md)
    bullets = extract_body_bullets(md)
    notes = extract_speaker_notes(md)
    n_mermaid = count_mermaid_blocks(md)
    inserts = extract_insert_markers(md)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Header strip ──
    header = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.7))
    hf = header.text_frame
    hf.word_wrap = True
    hp = hf.paragraphs[0]
    hr = hp.add_run()
    hr.text = f"Slide {slide_num} — {title_text}"
    hr.font.size = Pt(26)
    hr.font.bold = True
    hr.font.color.rgb = DEEP

    # ── Subtitle (if any) ──
    body_top = Inches(1.05)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.5))
        sf = sub.text_frame
        sf.word_wrap = True
        sp = sf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(14)
        sr.font.italic = True
        sr.font.color.rgb = SLATE
        body_top = Inches(1.45)

    # ── Diagram placeholder block (if any mermaid in MD) ──
    placeholder_lines: list[str] = []
    if n_mermaid > 0:
        placeholder_lines.append(
            f"📌 DIAGRAM PLACEHOLDER — {n_mermaid} Mermaid diagram"
            + ("s" if n_mermaid > 1 else "")
            + " in source MD."
        )
        placeholder_lines.append(
            "    Presenter: screenshot rendered diagram and replace this box."
        )
    if inserts:
        for ins in inserts[:3]:  # show up to 3 markers
            placeholder_lines.append(f"📌 {ins.strip()}")

    if placeholder_lines:
        ph = slide.shapes.add_textbox(Inches(0.4), body_top, Inches(12.5), Inches(1.6))
        pf = ph.text_frame
        pf.word_wrap = True
        pf.margin_top = Pt(8)
        pf.margin_bottom = Pt(8)
        pf.margin_left = Pt(10)
        pf.margin_right = Pt(10)
        for idx, line in enumerate(placeholder_lines):
            par = pf.paragraphs[0] if idx == 0 else pf.add_paragraph()
            run = par.add_run()
            run.text = line
            run.font.size = Pt(12)
            run.font.color.rgb = ACCENT
            run.font.bold = idx == 0
        # Frame fill
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(0xFD, 0xF4, 0xE3)
        ph.line.color.rgb = ACCENT
        body_top = body_top + Inches(1.75)

    # ── Body bullets ──
    bbox = slide.shapes.add_textbox(Inches(0.4), body_top, Inches(12.5), Inches(7.5) - body_top - Inches(0.3))
    bf = bbox.text_frame
    bf.word_wrap = True
    bf.margin_top = Pt(4)
    bf.margin_left = Pt(8)

    max_bullets = 22  # keep slide readable
    shown = bullets[:max_bullets]
    truncated = len(bullets) - len(shown)

    for idx, raw in enumerate(shown):
        para = bf.paragraphs[0] if idx == 0 else bf.add_paragraph()
        run = para.add_run()
        run.text = raw
        run.font.size = Pt(11)
        run.font.color.rgb = SLATE
        # Indent based on leading whitespace
        leading = len(raw) - len(raw.lstrip())
        if leading >= 4:
            para.level = 1

    if truncated > 0:
        para = bf.add_paragraph()
        run = para.add_run()
        run.text = f"  … and {truncated} more line(s) in source MD"
        run.font.size = Pt(10)
        run.font.italic = True
        run.font.color.rgb = TEAL

    # ── Footer ──
    foot = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
    ff = foot.text_frame
    fp = ff.paragraphs[0]
    fp.alignment = PP_ALIGN.RIGHT
    fr = fp.add_run()
    fr.text = f"AquaShield v2 · Slide {slide_num} / 16"
    fr.font.size = Pt(9)
    fr.font.color.rgb = TEAL

    # ── Speaker notes ──
    if notes:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = notes


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs)
    for i, path in enumerate(SLIDE_FILES, start=1):
        if not path.exists():
            print(f"  skip slide {i}: {path.name} not found")
            continue
        make_content_slide(prs, i, path)
        print(f"  built slide {i:>2}: {path.name}")

    prs.save(OUT)
    print(f"\n✅ Wrote {OUT}  ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
